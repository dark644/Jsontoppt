import json
from pptx import Presentation
from pptx.util import Pt, Inches
from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.textinput import TextInput
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.filechooser import FileChooserListView
from kivy.core.window import Window

class PPTGeneratorApp(App):
    def create_ppt_from_json(self, json_data):
        try:
            slides_data = json.loads(json_data)  # Parse JSON string
        except json.JSONDecodeError as e:
            return f"Error parsing JSON: {e}"

        if not isinstance(slides_data, list):
            return "Error: JSON must be a list of slides."

        ppt = Presentation()

        for slide in slides_data:
            slide_layout = ppt.slide_layouts[1]
            ppt_slide = ppt.slides.add_slide(slide_layout)

            title = slide.get("title", "Untitled Slide")
            ppt_slide.shapes.title.text = title

            content_placeholder = ppt_slide.placeholders[1]
            content_placeholder.left = Inches(0.5)
            content_placeholder.top = Inches(1.5)
            content_placeholder.width = Inches(9)
            content_placeholder.height = Inches(5.5)

            content = slide.get("content", "")
            bullet_points = slide.get("bullet_points", [])

            text_frame = content_placeholder.text_frame
            text_frame.clear()

            combined_text = content
            if bullet_points:
                combined_text += "\n\n" + "\n".join(f"→ {point}" for point in bullet_points)

            max_font_size = 24
            min_font_size = 12
            font_size = max_font_size

            while font_size >= min_font_size:
                text_frame.clear()
                p = text_frame.add_paragraph()
                p.text = combined_text
                p.font.size = Pt(font_size)

                text_height = len(combined_text.split("\n")) * font_size * 1.2
                slide_height = content_placeholder.height.inches * 72

                if text_height <= slide_height:
                    break
                font_size -= 2

        output_file = "Generated_Presentation.pptx"
        ppt.save(output_file)
        return f"Presentation saved as {output_file}"

    def build(self):
        Window.size = (600, 800)  # Set the app window size

        self.layout = BoxLayout(orientation="vertical", padding=10, spacing=10)

        self.instructions = Label(
            text="Paste your JSON input below and click 'Generate Presentation'.",
            size_hint=(1, 0.1),
        )
        self.layout.add_widget(self.instructions)

        self.text_input = TextInput(hint_text="Paste JSON here...", size_hint=(1, 0.6), multiline=True)
        self.layout.add_widget(self.text_input)

        self.generate_button = Button(
            text="Generate Presentation",
            size_hint=(1, 0.1),
            on_press=self.generate_presentation,
        )
        self.layout.add_widget(self.generate_button)

        self.result_label = Label(
            text="Output will appear here.",
            size_hint=(1, 0.2),
            halign="center",
            valign="middle",
        )
        self.result_label.bind(size=self.result_label.setter('text_size'))
        self.layout.add_widget(self.result_label)

        return self.layout

    def generate_presentation(self, instance):
        json_input = self.text_input.text.strip()
        if not json_input:
            self.result_label.text = "Please provide JSON input."
            return

        result = self.create_ppt_from_json(json_input)
        self.result_label.text = result


if __name__ == "__main__":
    PPTGeneratorApp().run()